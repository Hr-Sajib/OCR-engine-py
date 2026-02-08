from fastapi import FastAPI, File, UploadFile
from fastapi.responses import JSONResponse
from typing import List
import io
import os
import pdfplumber
import docx
from pptx import Presentation
from PIL import Image
from paddleocr import PaddleOCR

# --- Initialize OCR ---
ocr = PaddleOCR(use_angle_cls=True, lang='en')  # CPU mode, can be GPU if available

# --- FastAPI app ---
app = FastAPI(title="OCR & File Extraction API")

# --- Helper functions ---

def get_file_type(file: UploadFile):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext == ".pdf":
        return "pdf"
    elif ext in [".docx", ".doc"]:
        return "docx"
    elif ext in [".pptx", ".ppt"]:
        return "pptx"
    elif ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
        return "image"
    else:
        return "unknown"

def extract_pdf(file_bytes):
    text = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text.append(page_text.strip())
    return "\n".join(text)

def extract_docx(file_bytes):
    text = []
    doc = docx.Document(io.BytesIO(file_bytes))
    for para in doc.paragraphs:
        if para.text.strip():
            text.append(para.text.strip())
    return "\n".join(text)

def extract_pptx(file_bytes):
    text = []
    prs = Presentation(io.BytesIO(file_bytes))
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text.strip())
    return "\n".join(text)

def extract_image(file_bytes):
    image = Image.open(io.BytesIO(file_bytes)).convert("RGB")
    temp_path = "/tmp/temp_image.png"
    image.save(temp_path)

    result = ocr.ocr(temp_path, cls=True)
    text_lines = []
    for line in result:
        for word_info in line:
            text_lines.append(word_info[1][0])
    return "\n".join(text_lines)

def extract_text(file: UploadFile):
    file_bytes = file.file.read()
    file_type = get_file_type(file)

    if file_type == "pdf":
        return extract_pdf(file_bytes)
    elif file_type == "docx":
        return extract_docx(file_bytes)
    elif file_type == "pptx":
        return extract_pptx(file_bytes)
    elif file_type == "image":
        return extract_image(file_bytes)
    else:
        return "[Unsupported file type]"

# --- API endpoint ---
@app.post("/extract")
async def extract(files: List[UploadFile] = File(...)):
    if not files:
        return JSONResponse(content={"success": False, "message": "No files uploaded", "data": ""})

    extracted_texts = []
    for idx, file in enumerate(files):
        try:
            content = extract_text(file)
            extracted_texts.append(f"--- Source {idx+1}: {file.filename} ---\n{content}")
        except Exception as e:
            extracted_texts.append(f"--- Source {idx+1}: {file.filename} ---\n[Error extracting content: {str(e)}]")

    return JSONResponse(content={"success": True, "message": "Files extracted successfully", "data": "\n\n".join(extracted_texts)})

# --- Run standalone ---
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8100, reload=True)
